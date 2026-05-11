"""
競合クリエイティブ調査レポート PPTX ジェネレーター
"""

import io
import re
import html as html_module
import urllib.request
import urllib.error
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import datetime

# ─── カラーパレット（ライト・シンプル）────────────
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)  # 白
C_DARK      = RGBColor(0x1F, 0x29, 0x37)  # ダークグレー（文字用）
C_ACCENT    = RGBColor(0xC2, 0x7B, 0x8A)  # ソフトローズ
C_ACCENT2   = RGBColor(0xE8, 0xA8, 0x7C)  # ウォームピーチ
C_GRAY      = RGBColor(0x6B, 0x72, 0x80)  # グレー
C_LIGHT     = RGBColor(0xF8, 0xF9, 0xFA)  # 薄グレー（カード背景）
C_BORDER    = RGBColor(0xE5, 0xE7, 0xEB)  # ライトボーダー
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN     = RGBColor(0x05, 0x9E, 0x62)
C_RED       = RGBColor(0xDC, 0x26, 0x26)
C_ORANGE    = RGBColor(0xEA, 0x58, 0x0C)

STANCE_COLOR = {
    "攻め":      RGBColor(0x05, 0x9E, 0x62),
    "守り":      RGBColor(0x94, 0xA3, 0xB8),
    "変化中":    RGBColor(0xEA, 0x58, 0x0C),
    "変化中→攻め": RGBColor(0xE8, 0xA8, 0x7C),
}

# ─── ブランドデータ ──────────────────────────────
REPORT_DATE = "2026-04-06"

BRANDS = {
    "国内大手": [
        {
            "name": "ORBIS",
            "summary": "「SKINCARE LOUNGE」SAKURA ART EVENT（〜4/19）継続中\nアプリ600万DLキャンペーン（3/20〜4/19）。ポアレスキーププライマー継続展開\n体験型施設×SNS拡散の複合設計",
            "axis": "機能（毛穴・UV）＋体験型施設＋アプリ経済圏",
            "stance": "攻め",
            "url": "https://www.orbis.co.jp/news/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/2061/699/2061-699-c56fab8df75d9f1e32232c61fd5d4a84-1800x1273.jpg",
        },
        {
            "name": "KANEBO",
            "summary": "4/1〜X抽選「美容化粧水＋朝夜クリームサンプル 1万名プレゼント」新展開\n4/10「ルージュスターヴァイブラント」プリズムカラー限定2色（ピンクパール・チェリーレッド）\n「混ざり合い、調和する」キービジュアルテーマ",
            "axis": "哲学的情緒コピー＋SNS大規模獲得（1万名）＋限定色",
            "stance": "攻め",
            "url": "https://prtimes.jp/main/html/rd/p/000000881.000016220.html",
            "image_url": None,
            "ogp_page": "https://global.kanebo.com/ja/lp/new_lineup",
        },
        {
            "name": "DECORTE",
            "summary": "「イドラクラリティ」薬用ブライトニング化粧液をジェンダーレス訴求（WWDJapan掲載）\nルージュデコルテ「ティント＆プランプ」3色・5,500円継続\nスキンシャドウ デザイニングパレット追加色展開中",
            "axis": "高質感・ジェンダーレス・薬用成分",
            "stance": "守り",
            "url": "https://www.decorte.com/",
            "image_url": None,
            "ogp_page": "https://www.decorte.com/site/s/pointmakeup_202601.aspx",
        },
        {
            "name": "IPSA",
            "summary": "創業40周年「Fullness of skin and mind with ART」テーマに格上げ\n「ザ・タイムR アクア e」4/23発売予告（ロングセラー進化版）\n「春の浸透実感」キャンペーン＋日本茶着想限定パレット発売予告",
            "axis": "カウンセリング＋周年記念（ART訴求）＋限定デザイン",
            "stance": "変化中",
            "url": "https://www.ipsa.co.jp/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/100029/78/100029-78-80bce36734ad95e4426830719cae5d97-2564x1737.jpg",
        },
        {
            "name": "ELIXIR",
            "summary": "石田ゆり子×吉瀬美智子の2ミューズ体制で40〜50代に照準\n「デーケアレボリューション Brightening」トーンアップUV（3/21）\n「ピュアな明るさ、あげていこう」コピー。市場18年連続No.1実績訴求",
            "axis": "2ミューズ＋UV機能＋市場No.1実績",
            "stance": "攻め",
            "url": "https://www.shiseido.co.jp/elixir/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/5794/3032/5794-3032-5aeed01b7e8e52e8deb9637684ab3732-3900x3900.jpg",
        },
        {
            "name": "MUJI",
            "summary": "薬用ブライトニングシリーズ（4月末発売）が春の次の主軸\n天然由来100%・携帯サイズ展開（化粧水590円〜）\n4ライン体制（敏感肌・エイジケア・クリアケア・ブライトニング）完成",
            "axis": "シンプル・薬用・天然由来・プライスバリュー",
            "stance": "守り",
            "url": "https://www.muji.com/jp/ja/store/cmdty/section/T00020",
            "image_url": None,
            "ogp_page": "https://www.muji.com/jp/ja/store/cmdty/section/S10910",
        },
    ],
    "D2C・インフルエンサー系": [
        {
            "name": "Fujiko",
            "summary": "「マルチイージーパレット インパルスピンク」（2/5・5色2,195円）継続\nTikTok毛穴コンテンツでバイラル継続\nニュアンスラップティント先行サポーター100名UGC仕込み",
            "axis": "機能（毛穴）＋TikTokバイラル＋バズネーミング",
            "stance": "攻め",
            "url": "https://www.instagram.com/fujiko_brand/",
            "image_url": None,
            "ogp_page": "https://fujikobrand.com/",
        },
        {
            "name": "ETVOS",
            "summary": "「薬用 UVホワイトニングクリアセラム」キャンペーン（3/18〜4/14）\nミネラルプレストチーク10年ぶりリニューアル（5/1・SPF20/PA++・4色）\n5/1「Sugared Twilight」夏コレクション予告継続",
            "axis": "ミネラル成分＋薬用UV＋季節コレクション",
            "stance": "攻め",
            "url": "https://etvos.com/",
            "image_url": None,
            "ogp_page": "https://etvos.com/ss/",
        },
        {
            "name": "b.glen",
            "summary": "春季セール（3月）実施済み。QuSome®成分継続訴求\n定期便最大46%OFF・365日返品保証でLTV型EC設計維持\n2026年春の大型新商品発表は確認できず",
            "axis": "成分テクノロジー（QuSome）＋定期便LTV＋返品保証",
            "stance": "守り",
            "url": "https://www.bglen.net/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/21449/56/21449-56-0755631fc3f6a448b90e6713b4c0c823-750x812.jpg",
        },
        {
            "name": "Dr.Ci:Labo",
            "summary": "春コフレ「薬用BB×金ゲル＋VC100ローションミニ＋新UVビューティー液ミニ＋ポーチ」\nVC100系＋医薬部外品訴求継続\nクリニカルポジション強化中",
            "axis": "クリニカル＋成分（VC・金）＋コフレ限定感",
            "stance": "変化中",
            "url": "https://www.ci-labo.com/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/2961/613/2961-613-6ade05525d4a9a3bf4a8944a3125e715-1909x2700.jpg",
        },
        {
            "name": "UZU",
            "summary": "「UZU Eye Cream 00」が春の新作として投入（目元ケア拡張）\n「MADE IN JAPAN・無添加」安全性訴求を一貫維持\n目元専門ブランドとしてのアーティスティック路線継続",
            "axis": "安全性＋MADE IN JAPAN＋アート＋目元専門",
            "stance": "守り",
            "url": "https://www.uzu.team/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/10744/122/10744-122-1e307ecef737e62a443655340e34d5c7-2640x1760.jpg",
        },
        {
            "name": "WHOMEE",
            "summary": "WinC（MUSCAT GROUP子会社）が販売権取得・リブランディング本格着手\nイガリシノブがクリエイティブディレクション継続確定\n「シングルラメシャドウ パールベージュ」（3/9）等新色投入中",
            "axis": "タレント人格（イガリシノブ）＋ブランド刷新",
            "stance": "変化中",
            "url": "https://www.instagram.com/whomeeigari/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/19306/263/19306-263-d9d23c95c22697ef02f5b348840c2f0a-3900x2600.jpg",
        },
        {
            "name": "OWNCODE (By ttt.)",
            "summary": "本田翼フルプロデュース「By ttt.」ブランド\nスキンケア新作（リポクリアトナー・クリアセラミドクリーム）を本田翼SNSで自然体発信\n楽天・Amazon・Qoo10に出店拡大でEC動線強化",
            "axis": "タレント人格（本田翼）＋生活提案＋EC拡大",
            "stance": "攻め",
            "url": "https://by-ttt.jp/",
            "image_url": None,
            "ogp_page": "https://by-ttt.jp/",
        },
    ],
    "韓国コスメ": [
        {
            "name": "COSRX",
            "summary": "「RX ザ・ビタミンC23セラム」リニューアル（香り改良・処方改善）発売\nQoo10メガ割（3/27〜）での大規模EC展開継続\nスネイル成分×ナイアシン訴求の定番ライン安定稼働",
            "axis": "成分（スネイル・VC・ナイアシン）＋EC強化＋改良リニューアル",
            "stance": "攻め",
            "url": "https://www.instagram.com/cosrx_jp/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/147749/11/147749-11-7645aa2714d847cc8ea3d7946e1a35a2-3900x3900.jpg",
        },
        {
            "name": "Laneige",
            "summary": "「Juice Pop Box Lip Tint」（ジュースボックス型パッケージ・8色・2,310円）新登場\nNEO Cushion「The Matt」「The Glow」新バリエーション追加\nBTS JIN起用継続",
            "axis": "K-POPタレント（JIN）＋パッケージデザイン訴求＋リップライン拡充",
            "stance": "攻め",
            "url": "https://www.laneige.com/jp/ja/index.html",
            "image_url": None,
            "ogp_page": "https://www.laneige.com/jp/ja/index.html",
        },
        {
            "name": "innisfree",
            "summary": "「グリーンティー PDRN アイ＆スポットセラム」（2/1）口コミ拡大継続\n4/29「グリーンティーシードセラムN」発売予告\nボディケア新ライン「Ile Number Line」（フレグランス軸）追加",
            "axis": "K-POP（MINGYU）＋クリーンビューティー＋成分（PDRN・グリーンティー）",
            "stance": "攻め",
            "url": "https://www.innisfree.jp/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/31376/183/31376-183-ba62eda9e09b6e4f11261ee425e64623-2160x2700.jpg",
        },
        {
            "name": "MISSHA",
            "summary": "「ビタシープラス ジェル洗顔」（3/8）VitaC×黄土×毛穴訴求\nUGCリポスト中心のファン参加型設計継続\nクッションファンデのコスパ訴求維持。新展開は限定的",
            "axis": "コスパ＋UGC＋成分（VitaC）",
            "stance": "守り",
            "url": "https://www.missha.co.jp/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/23609/281/23609-281-4364a99fc4e7cd654c1bb247c6145448-2025x2700.jpg",
        },
        {
            "name": "rom&nd",
            "summary": "「ジューシーフラッシュリップオイル」4/7（明日）発売予定\nミッフィーコラボ（ZO&FRIENDS）継続。ハート型チーク等春コレクション展開\nTikTok「デパコス風プチプラ」コンテンツが安定拡散",
            "axis": "季節コレクション＋コラボIP＋TikTokバイラル",
            "stance": "攻め",
            "url": "https://www.instagram.com/romand_jp/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/76505/296/76505-296-c905d69beed81d0897dd3b34c4545db1-3200x1601.jpg",
        },
        {
            "name": "CLIO",
            "summary": "「シュガーアフタヌーンコレクション」継続\nToy Storyコラボ（プロアイパレットキューブ2色）で新規IP展開\nTXT（TomorrowXTogether）起用＋日本ローカライズ色「ムードピーチパレット」",
            "axis": "K-POP（TXT）＋コラボIP（Toy Story）＋日本ローカライズ",
            "stance": "攻め",
            "url": "https://cliocosmetic.jp/",
            "image_url": "https://prcdn.freetls.fastly.net/release_image/56098/98/56098-98-553cf7f3308e12e330b333e41038fc91-1920x1280.png",
        },
    ],
}

INSIGHTS = [
    {
        "icon": "✅",
        "title": "取り込むべきトレンド",
        "color": C_GREEN,
        "items": [
            "「哲学コピー×SNS大規模獲得」二段構え（KANEBO 1万名サンプルモデルが実証）",
            "UGC収集のタイミング設計（発売前後にハッシュタグキャンペーンをセット仕込み）",
            "TikTok→ECプラットフォーム複合ファネル整備（Qoo10メガ割×TikTok認知の連動）",
            "季節限定コレクション×IP（キャラクター）コラボ（rom&ndミッフィー等。話題創出コスト低）",
            "薬用・医薬部外品の成分訴求強化（MUJI・ETVOS・Dr.Ci:Laboで信頼性担保として機能）",
        ],
    },
    {
        "icon": "❌",
        "title": "過密・差別化困難な領域",
        "color": C_RED,
        "items": [
            "UV春商戦への単独参入（21社が集中投入。単独UV訴求では露出競争に埋没）",
            "成分名の羅列訴求（VC・ナイアシン・レチノール。「あって当然」の水準に到達）",
            "K-POPタレント一点集中起用（飽和＋コスト高騰。中小ブランドには非現実的）",
            "X/SNS抽選プレゼントキャンペーン（KANEBO・IPSA・MISSHAが同期間集中実施で単価上昇）",
        ],
    },
    {
        "icon": "🎯",
        "title": "差別化チャンスのポジション",
        "color": C_ACCENT,
        "items": [
            "「中価格帯（2,000〜4,000円）×ライフスタイル情緒訴求」空白ゾーン（最優先・変化なし）",
            "EC×小売オムニ体験の可視化（EC専業D2Cにも高価格デパコスにもできない複合強み）",
            "「日本プライドコスメ」情緒訴求（K-Beauty台頭の今こそ日本製品質を積極発信）",
            "WHOMEEリブランディング空席を狙う（イガリシノブファン層への今が接触チャンス）",
        ],
    },
]

WATCHLIST = [
    ("rom&nd", "「ジューシーフラッシュリップオイル」4/7発売の初動SNS反響とEC売上（4/8以降）"),
    ("IPSA", "「ザ・タイムR アクア e」4/23発売に向けたクリエイティブ展開（4/15前後）"),
    ("WHOMEE", "リブランディング後の新ビジュアル・新コピー公開タイミング（随時）"),
    ("ETVOS", "「Sugared Twilight」5/1発売前予告コンテンツとTikTok事前拡散（4月下旬）"),
    ("KANEBO", "1万名サンプルキャンペーン（〜4/15）の反響とフォロワー増加数（4/16以降）"),
    ("innisfree", "「グリーンティーシードセラムN」4/29発売と新ボディライン日本展開規模"),
    ("By ttt.", "楽天・Amazon・Qoo10への出店拡大後のEC売上推移（4〜5月）"),
    ("CLIO", "Toy Storyコラボの日本発売日確定と日本ローカライズ色展開（5月初旬）"),
]


# ─── ヘルパー関数 ────────────────────────────────

def fetch_image(url):
    """URLから画像を取得してBytesIOで返す。失敗時はNone"""
    if not url:
        return None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = resp.read()
        return io.BytesIO(data)
    except Exception:
        return None


def fetch_ogp_image(page_url):
    """WebページのOGP画像を取得してBytesIOで返す。失敗時はNone"""
    if not page_url:
        return None
    try:
        req = urllib.request.Request(page_url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=8) as r:
            raw = r.read().decode("utf-8", errors="ignore")
        m = re.search(r'<meta property="og:image"\s+content="([^"]+)"', raw)
        if not m:
            m = re.search(r'<meta content="([^"]+)"\s+property="og:image"', raw)
        if not m:
            return None
        img_url = html_module.unescape(m.group(1))
        return fetch_image(img_url)
    except Exception:
        return None


def set_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=12, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_hyperlink_text(slide, text, url, l, t, w, h, size=10, color=C_ACCENT, align=PP_ALIGN.LEFT):
    """クリッカブルリンク付きテキスト"""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.underline = True
    # スライドパートにリレーションシップを登録してrIdを取得
    rId = slide.part.relate_to(url, RT.HYPERLINK, is_external=True)
    rPr = run._r.get_or_add_rPr()
    hlinkClick = rPr.get_or_add_hlinkClick()
    hlinkClick.set(qn('r:id'), rId)
    return txBox


def stance_badge(slide, stance, l, t):
    """スタンスバッジを描画"""
    color = STANCE_COLOR.get(stance, C_GRAY)
    w, h = 1.2, 0.28
    shape = add_rect(slide, l, t, w, h, color)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = stance
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    return shape


# ─── スライド生成関数 ─────────────────────────────

def add_slide_header(slide, title):
    """全スライド共通のライトヘッダー（左アクセントバー＋タイトル＋下線）"""
    add_rect(slide, 0, 0, 0.1, 0.72, C_ACCENT)
    add_text(slide, title,
             0.28, 0.1, 12.8, 0.52, size=18, bold=True, color=C_DARK)
    add_rect(slide, 0, 0.72, 13.333, 0.02, C_BORDER)


def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # 左アクセントバー（縦長）
    add_rect(slide, 0, 0, 0.5, 7.5, C_ACCENT)

    # 右ページにコンテンツ
    # タグライン
    add_text(slide, "COMPETITOR CREATIVE REPORT",
             0.85, 1.6, 11.5, 0.45,
             size=11, bold=False, color=C_GRAY)

    # メインタイトル
    add_text(slide, "競合クリエイティブ\n調査レポート",
             0.85, 2.1, 11.5, 2.0,
             size=38, bold=True, color=C_DARK)

    # アクセントライン
    add_rect(slide, 0.85, 4.2, 3.5, 0.06, C_ACCENT)

    # 日付・概要
    add_text(slide, f"調査実施日　{REPORT_DATE}",
             0.85, 4.45, 8.0, 0.4,
             size=13, color=C_DARK)

    total = sum(len(v) for v in BRANDS.values())
    add_text(slide, f"調査対象 {total}ブランド ／ 国内大手・D2C系・韓国コスメ",
             0.85, 4.9, 11.0, 0.35,
             size=12, color=C_GRAY)

    add_text(slide, "2026年3月　春季クリエイティブ動向分析",
             0.85, 5.3, 11.0, 0.35,
             size=12, color=C_GRAY)

    # 下部ライン
    add_rect(slide, 0.5, 7.3, 12.833, 0.02, C_BORDER)


def make_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_header(slide, "エグゼクティブサマリー")

    findings = [
        ("01", "春商戦クライマックスで飽和ピーク",
         "KANEBO・IPSA・MISSHAが同期間に\nSNS抽選施策を集中実施。成分名訴求は\n21社中14社以上が実施し差別化力ゼロに。"),
        ("02", "WHOMEE・IPSA・COSRXが変化",
         "WHOMEEがリブランディング本格着手。\nIPSAが40周年テーマに格上げ。\nCOSRXがVC美容液を香り改良でリニューアル。"),
        ("03", "中価格×情緒の空白は依然無主",
         "2,000〜4,000円帯×ライフスタイル\n情緒訴求が全21社通じた空白ゾーン。\n自社の優先ポジションとして引き続き有望。"),
    ]

    for i, (num, title, body) in enumerate(findings):
        x = 0.3 + i * 4.3
        # カードBG（ライトグレー）
        add_rect(slide, x, 0.95, 4.1, 5.9, C_LIGHT)
        # 上部アクセントバー
        add_rect(slide, x, 0.95, 4.1, 0.06, C_ACCENT)

        # 番号
        add_text(slide, num, x + 0.15, 1.1, 0.7, 0.55,
                 size=22, bold=True, color=C_ACCENT)

        # タイトル
        add_text(slide, title, x + 0.15, 1.72, 3.8, 0.55,
                 size=14, bold=True, color=C_DARK)

        # 区切り線
        add_rect(slide, x + 0.15, 2.35, 3.6, 0.02, C_BORDER)

        # 本文
        add_text(slide, body, x + 0.15, 2.5, 3.8, 3.1,
                 size=11.5, color=C_GRAY)

    # フッター
    add_rect(slide, 0, 6.95, 13.333, 0.55, C_LIGHT)
    add_rect(slide, 0, 6.95, 0.1, 0.55, C_ACCENT2)
    add_text(slide, "全カテゴリ共通：Instagram Reels・TikTok縦型動画シフト と マルチユース商品訴求 が加速",
             0.25, 6.97, 12.9, 0.5, size=11, color=C_DARK, align=PP_ALIGN.CENTER)


def make_trend_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_header(slide, "今期のクリエイティブトレンド")

    # ── 左列：ビジュアルスタイル ──
    add_text(slide, "支配的ビジュアルスタイル", 0.3, 0.9, 6.0, 0.4,
             size=13, bold=True, color=C_DARK)

    styles = [
        ("クリーン成分クローズアップ", "ORBIS・IPSA・ELIXIR\n白背景＋テクスチャ寄り撮影"),
        ("シーズンテーマ統一", "ETVOS・Fujiko\n春夏でトーンを統一したコレクション感"),
        ("推し文化・ポートレート重視", "韓国コスメ全般\n人物を大きく配置、発色と共存"),
        ("縦型ライフスタイル映像", "全カテゴリ\nReels/TikTok向け縦型動画が急増"),
    ]
    for i, (title, desc) in enumerate(styles):
        y = 1.4 + i * 1.35
        add_rect(slide, 0.3, y, 0.06, 0.9, C_ACCENT)
        add_text(slide, title, 0.5, y, 5.6, 0.35, size=11, bold=True, color=C_DARK)
        add_text(slide, desc, 0.5, y + 0.35, 5.6, 0.6, size=9.5, color=C_GRAY)

    # 縦区切り
    add_rect(slide, 6.65, 0.85, 0.04, 6.2, RGBColor(0xDD, 0xDD, 0xDD))

    # ── 右列：訴求軸分布 ──
    add_text(slide, "主要訴求軸の分布", 6.85, 0.9, 6.0, 0.4,
             size=13, bold=True, color=C_DARK)

    axes = [
        ("成分・テクノロジー", 7, C_ACCENT),
        ("タレント・セレブリティ", 6, C_DARK),
        ("カラー・トレンド感", 5, RGBColor(0x8E, 0x44, 0xAD)),
        ("UGC・ユーザー参加型", 4, C_GREEN),
        ("ライフスタイル・哲学", 2, C_ORANGE),
        ("機能・マルチユース", 3, RGBColor(0x16, 0xA0, 0x85)),
    ]
    max_brands = 10
    for i, (label, count, color) in enumerate(axes):
        y = 1.4 + i * 0.85
        add_text(slide, label, 6.85, y, 3.0, 0.35, size=10, color=C_DARK)
        # バーグラフ
        bar_w = (count / max_brands) * 3.0
        add_rect(slide, 10.0, y + 0.04, bar_w, 0.28, color)
        add_text(slide, f"{count}ブランド", 10.0 + bar_w + 0.08, y + 0.04, 1.0, 0.28,
                 size=9, color=C_GRAY)

    # SNS傾向メモ
    add_rect(slide, 6.85, 6.6, 6.1, 0.82, C_LIGHT)
    add_rect(slide, 6.85, 6.6, 0.08, 0.82, C_ACCENT2)
    add_text(slide, "SNS vs 従来広告\nTV CM（大手中心）＋ Instagram全社＋ TikTok（韓国・D2C）\nリアル接点（駅広告・体験イベント）との融合が加速",
             6.98, 6.62, 5.9, 0.78, size=9, color=C_DARK)


def make_brand_slide(prs, category, brands):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_header(slide, f"ブランド別動向 ／ {category}")

    cols = 2 if len(brands) <= 4 else 3
    col_w = 12.8 / cols
    rows = (len(brands) + cols - 1) // cols
    card_h = (7.5 - 0.95 - 0.1) / rows - 0.12

    for i, brand in enumerate(brands):
        col = i % cols
        row = i // cols
        x = 0.27 + col * col_w
        y = 0.95 + row * (card_h + 0.12)

        # ── OGP画像取得（複数ソースを順番に試す）──────
        img_bytes = fetch_image(brand.get("image_url"))
        if not img_bytes:
            img_bytes = fetch_ogp_image(brand.get("ogp_page"))
        if not img_bytes:
            img_bytes = fetch_ogp_image(brand.get("article"))

        # ── 画像サイズ計算（アスペクト比保持） ────────
        MAX_IMG_W = col_w - 0.15
        MAX_IMG_H = 1.8             # ← 1.25→1.8 に拡大（よりビジュアル重視）
        img_h_actual = 0.0

        if img_bytes:
            try:
                img_bytes.seek(0)
                pil_img = PILImage.open(img_bytes)
                px_w, px_h = pil_img.size
                aspect = px_w / px_h

                h_if_full_w = MAX_IMG_W / aspect
                if h_if_full_w <= MAX_IMG_H:
                    final_img_w = MAX_IMG_W
                    final_img_h = h_if_full_w
                else:
                    final_img_h = MAX_IMG_H
                    final_img_w = MAX_IMG_H * aspect

                img_h_actual = final_img_h
                img_x = x + (MAX_IMG_W - final_img_w) / 2
                img_bytes.seek(0)

            except Exception as e:
                print(f"    ⚠️  {brand['name']} 画像サイズ取得失敗: {e}")
                img_bytes = None

        has_img = img_bytes is not None and img_h_actual > 0

        LINK_H  = 0.26
        txt_top = y + (img_h_actual + 0.06 if has_img else 0.06)

        # ── カード描画（ライトグレー背景＋上部スタンスバー）──
        add_rect(slide, x, y, col_w - 0.1, card_h, C_WHITE)
        # 外枠（ボーダー）
        add_rect(slide, x, y, col_w - 0.1, 0.02, C_BORDER)
        add_rect(slide, x, y + card_h - 0.02, col_w - 0.1, 0.02, C_BORDER)
        # スタンス色のトップバー
        stance_col = STANCE_COLOR.get(brand["stance"], C_GRAY)
        add_rect(slide, x, y, col_w - 0.1, 0.07, stance_col)

        IMG_TOP_GAP = 0.13  # 帯（0.07"）との間に小さなギャップを確保
        # 画像埋め込み
        if has_img:
            try:
                slide.shapes.add_picture(
                    img_bytes,
                    Inches(img_x), Inches(y + IMG_TOP_GAP),
                    Inches(final_img_w), Inches(final_img_h)
                )
                img_h_actual = final_img_h
                txt_top = y + IMG_TOP_GAP + img_h_actual + 0.06
                print(f"    🖼  {brand['name']}  ({final_img_w:.2f}\" × {final_img_h:.2f}\")")
            except Exception as e:
                print(f"    ⚠️  {brand['name']} 画像配置失敗: {e}")
                has_img = False
                txt_top = y + 0.1
        else:
            print(f"    ──  {brand['name']} 画像なし")
            txt_top = y + 0.1

        # ブランド名
        add_text(slide, brand["name"],
                 x + 0.12, txt_top, col_w - 0.55, 0.38,
                 size=12, bold=True, color=C_DARK)

        # スタンスバッジ
        stance_badge(slide, brand["stance"], x + col_w - 1.35, txt_top + 0.04)

        # 訴求軸（折り返しなし・カード幅いっぱい）
        add_text(slide, f"訴求：{brand['axis']}",
                 x + 0.12, txt_top + 0.38, col_w - 0.25, 0.26,
                 size=8, color=C_ACCENT, bold=True, wrap=False)

        # 区切り線
        add_rect(slide, x + 0.1, txt_top + 0.64, col_w - 0.3, 0.02, C_BORDER)

        # サマリー
        LINK_H   = 0.28
        link_y   = y + card_h - LINK_H - 0.08
        summary_h = max(link_y - (txt_top + 0.70) - 0.06, 0.25)
        add_text(slide, brand["summary"],
                 x + 0.12, txt_top + 0.70, col_w - 0.25, summary_h,
                 size=8.5, color=C_GRAY)

        # リンク（カード右端に右揃え）
        # カード右端 = x + col_w - 0.1（スタンスバー右端と同じ）
        CARD_RIGHT = x + col_w - 0.1
        LINK_W     = 1.1   # "Instagram" / "詳細リンク" が収まる幅
        link_label = "Instagram" if "instagram.com" in brand["url"] else "詳細リンク"
        if brand.get("youtube"):
            add_hyperlink_text(slide, "YouTube",
                               brand["youtube"],
                               CARD_RIGHT - LINK_W * 2 - 0.1, link_y, LINK_W, LINK_H,
                               size=8, align=PP_ALIGN.RIGHT)
            add_hyperlink_text(slide, link_label,
                               brand["url"],
                               CARD_RIGHT - LINK_W, link_y, LINK_W, LINK_H,
                               size=8, align=PP_ALIGN.RIGHT)
        elif brand.get("article"):
            add_hyperlink_text(slide, link_label,
                               brand["url"],
                               CARD_RIGHT - LINK_W * 2 - 0.1, link_y, LINK_W, LINK_H,
                               size=8, align=PP_ALIGN.RIGHT)
            add_hyperlink_text(slide, "記事",
                               brand["article"],
                               CARD_RIGHT - LINK_W, link_y, LINK_W, LINK_H,
                               size=8, align=PP_ALIGN.RIGHT)
        else:
            add_hyperlink_text(slide, link_label,
                               brand["url"],
                               CARD_RIGHT - LINK_W, link_y, LINK_W, LINK_H,
                               size=8, align=PP_ALIGN.RIGHT)


def make_competitive_map(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_header(slide, "競合ポジショニングマップ")

    # 4象限のBG
    mx, my = 0.6, 1.0
    mw, mh = 12.1, 5.8
    half_w, half_h = mw / 2, mh / 2

    # 象限色
    add_rect(slide, mx,          my,          half_w, half_h, RGBColor(0xF5, 0xE8, 0xEE))  # 高×情緒
    add_rect(slide, mx + half_w, my,          half_w, half_h, RGBColor(0xE8, 0xF0, 0xF8))  # 高×機能
    add_rect(slide, mx,          my + half_h, half_w, half_h, RGBColor(0xFD, 0xF3, 0xE8))  # 低×情緒
    add_rect(slide, mx + half_w, my + half_h, half_w, half_h, RGBColor(0xE8, 0xF5, 0xED))  # 低×機能

    # 中央十字線（先に描画してテキストが前面に来るようにする）
    add_rect(slide, mx + half_w - 0.02, my, 0.04, mh, RGBColor(0xCC, 0xCC, 0xCC))
    add_rect(slide, mx, my + half_h - 0.02, mw, 0.04, RGBColor(0xCC, 0xCC, 0xCC))

    # 軸ラベル（十字線より後に描画 → 最前面）
    add_text(slide, "高価格帯", mx, my - 0.05, mw, 0.35,
             size=10, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "低価格帯", mx, my + mh, mw, 0.35,
             size=10, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "← 情緒訴求", mx, my + half_h - 0.18, half_w, 0.35,
             size=10, bold=True, color=RGBColor(0xB5, 0x60, 0x7A), align=PP_ALIGN.CENTER)
    add_text(slide, "機能訴求 →", mx + half_w, my + half_h - 0.18, half_w, 0.35,
             size=10, bold=True, color=RGBColor(0x2E, 0x6D, 0xA4), align=PP_ALIGN.CENTER)

    # 象限タイトル
    add_text(slide, "高価格×情緒", mx + 0.1, my + 0.1, half_w - 0.2, 0.3,
             size=9, bold=True, color=C_ACCENT)
    add_text(slide, "高価格×機能", mx + half_w + 0.1, my + 0.1, half_w - 0.2, 0.3,
             size=9, bold=True, color=RGBColor(0x2E, 0x6D, 0xA4))
    add_text(slide, "低価格×情緒  ⚠️ 過密", mx + 0.1, my + half_h + 0.1, half_w - 0.2, 0.3,
             size=9, bold=True, color=C_ORANGE)
    add_text(slide, "低価格×機能", mx + half_w + 0.1, my + half_h + 0.1, half_w - 0.2, 0.3,
             size=9, bold=True, color=C_GREEN)

    # ブランド配置データ
    brand_positions = {
        # (x比率0-1, y比率0-1) ※x=0が左(情緒), x=1が右(機能), y=0が上(高価格), y=1が下(低価格)
        "KANEBO":   (0.15, 0.2),
        "DECORTE":  (0.3, 0.15),
        "ELIXIR":   (0.55, 0.25),
        "IPSA":     (0.75, 0.2),
        "b.glen":   (0.85, 0.3),
        "Dr.Ci:Labo": (0.8, 0.4),
        "ORBIS":    (0.7, 0.65),
        "COSRX":    (0.82, 0.72),
        "Laneige":  (0.2, 0.6),
        "rom&nd":   (0.25, 0.78),
        "CLIO":     (0.3, 0.7),
        "ETVOS":    (0.18, 0.72),
        "Fujiko":   (0.6, 0.75),
        "UZU":      (0.22, 0.85),
        "WHOMEE":   (0.15, 0.9),
        "ELIXIR (Urfle)": (0.65, 0.7),
        "innisfree": (0.38, 0.58),
    }

    for name, (rx, ry) in brand_positions.items():
        bx = mx + rx * mw - 0.5
        by = my + ry * mh - 0.15
        # 小さなラベル
        add_text(slide, name, bx, by, 1.1, 0.28, size=7.5, color=C_DARK)


def make_insights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_header(slide, "自社への示唆・アクション案")

    col_w = 4.2
    for i, insight in enumerate(INSIGHTS):
        x = 0.3 + i * 4.4
        y = 0.95

        # 列カード（ライトグレー背景）
        add_rect(slide, x, y, col_w, 6.38, C_LIGHT)
        # 上部カラーバー
        add_rect(slide, x, y, col_w, 0.08, insight["color"])

        # 列ヘッダー
        add_rect(slide, x, y + 0.08, col_w, 0.65, C_WHITE)
        add_text(slide, f"{insight['icon']} {insight['title']}",
                 x + 0.15, y + 0.12, col_w - 0.25, 0.55,
                 size=11, bold=True, color=C_DARK)

        # アイテムリスト
        for j, item in enumerate(insight["items"]):
            iy = y + 0.88 + j * 1.42
            # 左アクセント点
            add_rect(slide, x + 0.12, iy + 0.13, 0.08, 0.08, insight["color"])
            add_text(slide, item, x + 0.28, iy, col_w - 0.4, 1.28,
                     size=10, color=C_DARK)
            if j < len(insight["items"]) - 1:
                add_rect(slide, x + 0.1, iy + 1.32, col_w - 0.2, 0.02, C_BORDER)


def make_watchlist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_slide_header(slide, "次回調査までのウォッチリスト")

    # テーブルヘッダー
    add_rect(slide, 0.3, 0.9, 1.9, 0.42, C_DARK)
    add_text(slide, "ブランド", 0.3, 0.9, 1.9, 0.42,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.25, 0.9, 10.8, 0.42, C_DARK)
    add_text(slide, "注目ポイント", 2.25, 0.9, 10.8, 0.42,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for i, (brand, point) in enumerate(WATCHLIST):
        y = 1.4 + i * 0.9
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        add_rect(slide, 0.3, y, 1.9, 0.82, bg)
        add_rect(slide, 2.25, y, 10.8, 0.82, bg)
        # 左アクセント
        add_rect(slide, 0.3, y, 0.06, 0.82, C_ACCENT2)
        add_text(slide, brand, 0.44, y + 0.1, 1.7, 0.6,
                 size=11, bold=True, color=C_DARK)
        add_text(slide, point, 2.35, y + 0.09, 10.6, 0.7,
                 size=10, color=C_GRAY)

    # フッター
    add_rect(slide, 0, 7.32, 13.333, 0.18, C_LIGHT)
    add_rect(slide, 0, 7.32, 13.333, 0.02, C_BORDER)
    add_text(slide, f"調査実施日：{REPORT_DATE}　　次回調査：2026-04-01",
             0.3, 7.33, 12.7, 0.16,
             size=8.5, color=C_GRAY, align=PP_ALIGN.CENTER)


# ─── メイン ──────────────────────────────────────

def build_pptx(output_path):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("📊 スライド生成中...")
    make_cover(prs)
    print("  ✅ Cover")

    make_executive_summary(prs)
    print("  ✅ Executive Summary")

    make_trend_overview(prs)
    print("  ✅ Trend Overview")

    for category, brands in BRANDS.items():
        make_brand_slide(prs, category, brands)
        print(f"  ✅ {category}")

    make_competitive_map(prs)
    print("  ✅ Competitive Map")

    make_insights(prs)
    print("  ✅ Insights")

    make_watchlist(prs)
    print("  ✅ Watchlist")

    prs.save(output_path)
    print(f"\n✨ 保存完了: {output_path}")


if __name__ == "__main__":
    out = f"/Users/dz0019/Desktop/claude_management/output/trends/{REPORT_DATE}-competitor-creative-report.pptx"
    build_pptx(out)
